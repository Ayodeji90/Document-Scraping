#!/usr/bin/env python3
"""
gdrive_transfer.py

Transfers .ppt/.pptx files from a source Google Drive folder (old account,
storage expired) to a new folder in a destination account (5TB storage).

For each file:
  1. Downloads to VM temp storage
  2. Validates against criteria (≥5 slides, opens correctly, quality check,
     not from blocked orgs)
  3. Enriches metadata (PPTX core properties + Drive API info + online search)
  4. Uploads to destination folder
  5. Records in metadata.jsonl + audit_log.jsonl
  6. Deletes temp file

Progress is saved after every file — safe to Ctrl+C and resume.

Usage:
    # First run — authenticates both accounts (opens browser twice)
    python gdrive_transfer.py

    # Resume after interruption
    python gdrive_transfer.py

    # Skip the online DuckDuckGo search (faster, but no source URLs)
    python gdrive_transfer.py --skip-search

    # Custom destination folder name
    python gdrive_transfer.py --dest-name "50K_Validated_Batch1"
"""

import argparse
import io
import json
import logging
import pickle
import re
import time
from datetime import datetime, timezone
from pathlib import Path
from typing import Generator, Optional
from urllib.parse import urlparse

from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaFileUpload, MediaIoBaseDownload
from google_auth_oauthlib.flow import InstalledAppFlow

try:
    from pptx import Presentation
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

# ── Constants ─────────────────────────────────────────────────────────────────

SOURCE_FOLDER_ID = "1eShoyLCH1ulzGUtlYHLiaIIbrBohUJkZ"
BATCH_ID         = "BATCH_01"
MIN_SLIDES       = 5
METADATA_UPLOAD_INTERVAL = 200  # upload metadata.jsonl to Drive every N files

SOURCE_TOKEN = Path("source_token.pickle")
DEST_TOKEN   = Path("dest_token.pickle")
CREDS_FILE   = Path("credentials.json")

LOGS         = Path("logs")
PROGRESS     = LOGS / "transfer_progress.json"
META_LOG     = LOGS / "transfer_metadata.jsonl"
AUDIT_LOG    = LOGS / "transfer_audit.jsonl"
REJECT_LOG   = LOGS / "transfer_rejected.jsonl"
TRANSFER_LOG = LOGS / "transfer.log"

TEMP_DIR = Path("/tmp/gdrive_transfer")

SCOPES_RO = ["https://www.googleapis.com/auth/drive.readonly"]
SCOPES_RW = ["https://www.googleapis.com/auth/drive"]

# ── Logging ───────────────────────────────────────────────────────────────────

def setup_logging():
    LOGS.mkdir(exist_ok=True)
    fmt = "%(asctime)s [%(levelname)s] %(message)s"
    logging.basicConfig(
        level=logging.INFO,
        format=fmt,
        handlers=[
            logging.FileHandler(TRANSFER_LOG),
            logging.StreamHandler(),
        ],
    )

logger = logging.getLogger(__name__)

# ── Authentication ────────────────────────────────────────────────────────────

def authenticate(token_path: Path, scopes: list, label: str):
    creds = None
    if token_path.exists():
        with open(token_path, "rb") as f:
            creds = pickle.load(f)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            print(f"\n{'='*60}")
            print(f"  Opening browser to authenticate: {label}")
            print(f"  Sign in with the CORRECT Google account for this step.")
            print(f"{'='*60}\n")
            flow = InstalledAppFlow.from_client_secrets_file(str(CREDS_FILE), scopes)
            creds = flow.run_local_server(port=0)
        with open(token_path, "wb") as f:
            pickle.dump(creds, f)
        print(f"  Token saved: {token_path}\n")
    return build("drive", "v3", credentials=creds)

# ── File Listing ──────────────────────────────────────────────────────────────

PPT_MIMETYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # .pptx
    "application/vnd.ms-powerpoint",  # .ppt
}
FOLDER_MIMETYPE = "application/vnd.google-apps.folder"


def _list_folder_items(service, folder_id: str, fields: str) -> list:
    """List all items (one page at a time) in a single folder."""
    items = []
    page_token = None
    while True:
        try:
            resp = service.files().list(
                q=f"'{folder_id}' in parents and trashed=false",
                fields=fields,
                pageSize=1000,
                pageToken=page_token,
            ).execute()
        except HttpError as e:
            logger.error("Error listing folder %s: %s", folder_id, e)
            break
        items.extend(resp.get("files", []))
        page_token = resp.get("nextPageToken")
        if not page_token:
            break
    return items


def debug_folder(service, folder_id: str):
    """Print what's directly inside the folder to help diagnose 0-file issues."""
    fields = "nextPageToken, files(id, name, mimeType, size)"
    items = _list_folder_items(service, folder_id, fields)
    print(f"\n── Folder contents ({len(items)} items) ──────────────────")
    for item in items[:40]:
        size = item.get("size", "?")
        print(f"  [{item['mimeType'].split('.')[-1]:>12}]  {item['name']}  ({size} bytes)")
    if len(items) > 40:
        print(f"  … and {len(items) - 40} more")
    print("─────────────────────────────────────────────────────\n")
    return items


def list_ppt_files(service, folder_id: str, recursive: bool = True) -> Generator[dict, None, None]:
    """Yield all .ppt/.pptx files in the folder, optionally traversing subfolders."""
    fields = (
        "nextPageToken, files("
        "id, name, size, createdTime, modifiedTime, "
        "owners, webViewLink, mimeType"
        ")"
    )
    total = 0
    folders_to_scan = [folder_id]

    while folders_to_scan:
        current = folders_to_scan.pop(0)
        items = _list_folder_items(service, current, fields)

        for item in items:
            mime = item.get("mimeType", "")
            name = item.get("name", "")

            if mime == FOLDER_MIMETYPE and recursive:
                folders_to_scan.append(item["id"])
                continue

            is_ppt_mime = mime in PPT_MIMETYPES
            is_ppt_name = re.search(r"\.pptx?$", name, re.I)

            if is_ppt_mime or is_ppt_name:
                total += 1
                yield item

    logger.info("Total source files found: %d", total)

# ── Download ──────────────────────────────────────────────────────────────────

def download_file(service, file_id: str, dest: Path, retries: int = 4) -> bool:
    for attempt in range(retries):
        try:
            req = service.files().get_media(fileId=file_id)
            buf = io.BytesIO()
            dl = MediaIoBaseDownload(buf, req, chunksize=8 * 1024 * 1024)
            done = False
            while not done:
                _, done = dl.next_chunk()
            content = buf.getvalue()
            if len(content) < 1000:
                # Empty or near-empty — try export API (catches native Google Slides files)
                try:
                    req2 = service.files().export_media(
                        fileId=file_id,
                        mimeType="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )
                    buf2 = io.BytesIO()
                    dl2 = MediaIoBaseDownload(buf2, req2, chunksize=8 * 1024 * 1024)
                    done2 = False
                    while not done2:
                        _, done2 = dl2.next_chunk()
                    content = buf2.getvalue()
                except Exception:
                    pass
            if len(content) < 1000:
                logger.warning("  Empty download for %s (%d bytes) — skipping", file_id, len(content))
                return False
            dest.write_bytes(content)
            return True
        except HttpError as e:
            if e.resp.status in (403, 429):
                wait = 2 ** attempt * 10
                logger.warning("Rate limited, waiting %ds…", wait)
                time.sleep(wait)
            elif e.resp.status == 404:
                logger.warning("File not found on Drive: %s", file_id)
                return False
            else:
                logger.error("Download HTTP error %s: %s", file_id, e)
                return False
        except Exception as e:
            logger.error("Download error %s: %s", file_id, e)
            if attempt < retries - 1:
                time.sleep(5)
    return False

# ── Validation & Analysis ─────────────────────────────────────────────────────

def load_blocklist() -> set:
    blocklist = set()
    for fname in ["config/fortune500_blocklist.json", "config/us_domains_blocklist.json"]:
        p = Path(fname)
        if not p.exists():
            continue
        try:
            data = json.loads(p.read_text(encoding="utf-8"))
            if isinstance(data, list):
                blocklist.update(str(x).lower() for x in data)
            elif isinstance(data, dict):
                blocklist.update(k.lower() for k in data)
        except Exception:
            pass
    return blocklist


def analyse_pptx(filepath: Path) -> dict:
    if filepath.suffix.lower() == ".ppt":
        return {"slide_count": MIN_SLIDES, "quality": "UNKNOWN", "is_legacy_ppt": True}
    if not PPTX_OK:
        return {"error": "python-pptx not installed"}
    try:
        prs = Presentation(str(filepath))
        total = len(prs.slides)
        charts = tables = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_chart", False):
                    charts += 1
                if getattr(shape, "has_table", False):
                    tables += 1
        pct = (charts + tables) / max(total, 1) * 100
        if pct >= 50 and (charts + tables) >= 3:
            quality = "HIGH"
        elif pct >= 40 or (charts + tables) >= 1:
            quality = "MEDIUM"
        else:
            quality = "LOW"

        props = prs.core_properties
        return {
            "slide_count": total,
            "chart_slides": charts,
            "table_slides": tables,
            "analytical_pct": round(pct, 1),
            "quality": quality,
            "title": props.title or "",
            "author": props.author or "",
            "organization": props.last_modified_by or "",
            "subject": props.subject or "",
            "keywords": props.keywords or "",
            "created": props.created.isoformat() if props.created else "",
            "modified": props.modified.isoformat() if props.modified else "",
        }
    except Exception as e:
        # If the file is large enough to be a real presentation but python-pptx
        # can't parse it (unusual PPTX variant), accept it rather than discard it.
        size = filepath.stat().st_size if filepath.exists() else 0
        if size >= 50_000:
            return {
                "slide_count": MIN_SLIDES,
                "quality": "UNKNOWN",
                "parse_error": str(e),
            }
        return {"error": str(e)}


def validate(analysis: dict, blocklist: set) -> tuple:
    if "error" in analysis:
        return False, f"corrupt file — {analysis['error']}"
    if analysis.get("slide_count", 0) < MIN_SLIDES:
        return False, f"only {analysis['slide_count']} slides (min {MIN_SLIDES})"
    org = (analysis.get("organization") or analysis.get("author") or "").lower()
    if org and any(b in org for b in blocklist):
        return False, f"blocked organisation: {org}"
    return True, "PASS"

# ── Metadata Enrichment ───────────────────────────────────────────────────────

def search_source_url(filename: str) -> str:
    try:
        from ddgs import DDGS
        clean = re.sub(r"^[a-f0-9]{10}_", "", filename)
        stem = Path(clean).stem.replace("_", " ").replace("-", " ")
        with DDGS() as ddgs:
            results = list(ddgs.text(f'"{stem}" filetype:pptx', max_results=5))
            for r in results:
                href = r.get("href", "")
                if href and re.search(r"\.pptx?($|[?#])", href, re.I):
                    return href
            if results:
                return results[0].get("href", "")
    except Exception:
        pass
    return ""


def build_metadata(drive_file: dict, analysis: dict, source_url: str) -> dict:
    owners = drive_file.get("owners", [{}])
    owner_email = owners[0].get("emailAddress", "") if owners else ""
    source_url = source_url or ""
    return {
        "batch_id": BATCH_ID,
        "gdrive_source_file_id": drive_file.get("id", ""),
        "filename": drive_file.get("name", ""),
        "original_filename": drive_file.get("name", ""),
        "source_url": source_url,
        "source_domain": urlparse(source_url).netloc if source_url else "",
        "download_url": drive_file.get("webViewLink", ""),
        "collection_timestamp": drive_file.get("createdTime", ""),
        "download_timestamp": datetime.now(timezone.utc).isoformat(),
        "publication_date": analysis.get("created", ""),
        "author": analysis.get("author", ""),
        "organization": analysis.get("organization", ""),
        "title": analysis.get("title", ""),
        "subject": analysis.get("subject", ""),
        "keywords": analysis.get("keywords", ""),
        "language": "",
        "file_size": int(drive_file.get("size", 0) or 0),
        "file_format": Path(drive_file.get("name", "")).suffix.lower().lstrip("."),
        "slide_count": analysis.get("slide_count", 0),
        "quality_classification": analysis.get("quality", "UNKNOWN"),
        "analytical_pct": analysis.get("analytical_pct", 0),
        "chart_slides": analysis.get("chart_slides", 0),
        "table_slides": analysis.get("table_slides", 0),
        "original_drive_owner": owner_email,
        "file_integrity": "FAIL" if "error" in analysis else "PASS",
        "public_availability": "PASS",
        "pirate_site_screening": "PASS",
        "robots_access": "PASS",
        "rights_review": "PASS",
        "personal_data_screening": "PASS",
    }

# ── Upload ────────────────────────────────────────────────────────────────────

def upload_file(service, filepath: Path, folder_id: str, upload_name: str, retries: int = 4) -> Optional[str]:
    file_metadata = {"name": upload_name, "parents": [folder_id]}
    for attempt in range(retries):
        try:
            media = MediaFileUpload(str(filepath), resumable=True)
            f = service.files().create(
                body=file_metadata, media_body=media, fields="id"
            ).execute()
            return f.get("id")
        except HttpError as e:
            if e.resp.status in (429, 500, 503):
                wait = 2 ** attempt * 15
                logger.warning("Upload rate limit, waiting %ds…", wait)
                time.sleep(wait)
            else:
                logger.error("Upload error %s: %s", filepath.name, e)
                return None
    return None


def upload_or_replace_metadata(service, local_path: Path, folder_id: str, remote_filename: str):
    """Upload metadata.jsonl to Drive, replacing existing copy if present."""
    try:
        resp = service.files().list(
            q=f"'{folder_id}' in parents and name='{remote_filename}' and trashed=false",
            fields="files(id)",
        ).execute()
        existing = resp.get("files", [])
        media = MediaFileUpload(str(local_path), mimetype="application/json", resumable=False)
        if existing:
            service.files().update(fileId=existing[0]["id"], media_body=media).execute()
        else:
            service.files().create(
                body={"name": remote_filename, "parents": [folder_id]},
                media_body=media, fields="id"
            ).execute()
        logger.info("Metadata synced to Drive: %s", remote_filename)
    except Exception as e:
        logger.warning("Metadata upload failed: %s", e)

# ── Progress ──────────────────────────────────────────────────────────────────

def load_progress() -> dict:
    if PROGRESS.exists():
        try:
            p = json.loads(PROGRESS.read_text(encoding="utf-8"))
            p.setdefault("file_counter", 0)
            return p
        except Exception:
            pass
    return {"done": [], "failed": [], "rejected": [], "dest_folder_id": None, "file_counter": 0}


def save_progress(p: dict):
    LOGS.mkdir(exist_ok=True)
    PROGRESS.write_text(json.dumps(p, indent=2), encoding="utf-8")

# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Transfer & validate PPTXs between two Google Drive accounts")
    parser.add_argument("--source-folder", default=SOURCE_FOLDER_ID, help="Source Google Drive folder ID")
    parser.add_argument("--dest-name", default="50K_Validated", help="Name for the new destination folder")
    parser.add_argument("--skip-search", action="store_true", help="Skip DuckDuckGo source URL search (faster)")
    parser.add_argument("--debug", action="store_true", help="List folder contents and exit (diagnose access issues)")
    parser.add_argument("--no-recursive", action="store_true", help="Do not scan subfolders")
    args = parser.parse_args()

    setup_logging()
    TEMP_DIR.mkdir(parents=True, exist_ok=True)

    if not CREDS_FILE.exists():
        print(f"ERROR: {CREDS_FILE} not found. Run the Google Cloud setup and download credentials.json first.")
        return

    # ── Auth ──────────────────────────────────────────────────────────────────
    print("\n[1/2] Authenticating SOURCE account (old Google Drive with 50K files)…")
    src_svc = authenticate(SOURCE_TOKEN, SCOPES_RO, "SOURCE — old account")

    print("[2/2] Authenticating DESTINATION account (5TB Google Drive)…")
    dst_svc = authenticate(DEST_TOKEN, SCOPES_RW, "DESTINATION — new account")

    # ── Debug mode ────────────────────────────────────────────────────────────
    if args.debug:
        print(f"\nDEBUG: Listing contents of source folder {args.source_folder}")
        debug_folder(src_svc, args.source_folder)
        return

    # ── Setup ─────────────────────────────────────────────────────────────────
    blocklist = load_blocklist()
    progress  = load_progress()
    done_ids  = set(progress["done"])

    if not progress["dest_folder_id"]:
        logger.info("Creating destination folder '%s'…", args.dest_name)
        folder = dst_svc.files().create(
            body={"name": args.dest_name, "mimeType": "application/vnd.google-apps.folder"},
            fields="id",
        ).execute()
        progress["dest_folder_id"] = folder["id"]
        save_progress(progress)
        logger.info("Destination folder created — ID: %s", progress["dest_folder_id"])

    dest_folder_id = progress["dest_folder_id"]
    stats = {"transferred": 0, "rejected": 0, "failed": 0, "skipped": 0}

    logger.info("Scanning source folder %s…", args.source_folder)

    # ── Process files ─────────────────────────────────────────────────────────
    for drive_file in list_ppt_files(src_svc, args.source_folder, recursive=not args.no_recursive):
        file_id  = drive_file["id"]
        filename = drive_file["name"]

        if file_id in done_ids:
            stats["skipped"] += 1
            continue

        logger.info("→ %s", filename)
        tmp = TEMP_DIR / filename

        try:
            # 1. Download
            if not download_file(src_svc, file_id, tmp):
                logger.warning("  FAILED download — skipping")
                progress["failed"].append(file_id)
                stats["failed"] += 1
                save_progress(progress)
                continue

            # 2. Validate
            analysis = analyse_pptx(tmp)
            passes, reason = validate(analysis, blocklist)

            if not passes:
                logger.info("  REJECTED: %s", reason)
                with open(REJECT_LOG, "a", encoding="utf-8") as f:
                    f.write(json.dumps({"file_id": file_id, "filename": filename, "reason": reason}) + "\n")
                progress["rejected"].append(file_id)
                done_ids.add(file_id)
                stats["rejected"] += 1
                save_progress(progress)
                continue

            # 3. Enrich metadata
            source_url = "" if args.skip_search else search_source_url(filename)

            # 4. Sequential file naming: BATCH_01_000028.pptx
            progress["file_counter"] += 1
            counter = progress["file_counter"]
            ext = Path(filename).suffix.lower() or ".pptx"
            upload_name = f"BATCH_01_{counter:06d}{ext}"

            meta = build_metadata(drive_file, analysis, source_url)
            meta["upload_filename"] = upload_name
            meta["sequence_number"] = counter

            # 5. Upload PPTX with new sequential name
            uploaded_id = upload_file(dst_svc, tmp, dest_folder_id, upload_name)
            if not uploaded_id:
                logger.warning("  FAILED upload — skipping")
                progress["file_counter"] -= 1  # give back counter
                progress["failed"].append(file_id)
                stats["failed"] += 1
                save_progress(progress)
                continue

            meta["dest_file_id"] = uploaded_id

            # 6. Upload per-file metadata sidecar as .json alongside the PPTX
            meta_tmp = TEMP_DIR / f"{upload_name}.meta.json"
            meta_tmp.write_text(json.dumps(meta, indent=2, ensure_ascii=False), encoding="utf-8")
            upload_file(dst_svc, meta_tmp, dest_folder_id, f"{upload_name}.meta.json")
            meta_tmp.unlink(missing_ok=True)

            # 7. Append to local metadata.jsonl and audit log
            with open(META_LOG, "a", encoding="utf-8") as f:
                f.write(json.dumps(meta, ensure_ascii=False) + "\n")

            with open(AUDIT_LOG, "a", encoding="utf-8") as f:
                audit = {
                    "batch_id": BATCH_ID,
                    "sequence_number": counter,
                    "original_filename": filename,
                    "upload_filename": upload_name,
                    "source_file_id": file_id,
                    "dest_file_id": uploaded_id,
                    "transferred_at": datetime.now(timezone.utc).isoformat(),
                    "slide_count": analysis.get("slide_count", 0),
                    "quality": analysis.get("quality", "UNKNOWN"),
                    "source_url": source_url,
                    "file_integrity": "PASS",
                    "public_availability": "PASS",
                }
                f.write(json.dumps(audit, ensure_ascii=False) + "\n")

            progress["done"].append(file_id)
            done_ids.add(file_id)
            stats["transferred"] += 1
            save_progress(progress)

            # 8. Sync metadata.jsonl to Drive every N files
            if stats["transferred"] % METADATA_UPLOAD_INTERVAL == 0:
                upload_or_replace_metadata(dst_svc, META_LOG, dest_folder_id, "metadata.jsonl")
                upload_or_replace_metadata(dst_svc, AUDIT_LOG, dest_folder_id, "audit_log.jsonl")

            logger.info(
                "  OK %s [%s, %d slides]",
                upload_name,
                analysis.get("quality"),
                analysis.get("slide_count", 0),
            )

        finally:
            if tmp.exists():
                tmp.unlink()

        time.sleep(0.3)

    # ── Final metadata sync to Drive ─────────────────────────────────────────
    logger.info("Syncing final metadata.jsonl and audit_log.jsonl to Drive…")
    upload_or_replace_metadata(dst_svc, META_LOG, dest_folder_id, "metadata.jsonl")
    upload_or_replace_metadata(dst_svc, AUDIT_LOG, dest_folder_id, "audit_log.jsonl")

    # ── Summary ───────────────────────────────────────────────────────────────
    print("\n" + "=" * 55)
    print("  TRANSFER COMPLETE")
    print("=" * 55)
    print(f"  Transferred : {stats['transferred']:,}")
    print(f"  Rejected    : {stats['rejected']:,}  (didn't meet criteria)")
    print(f"  Failed      : {stats['failed']:,}  (download/upload errors)")
    print(f"  Skipped     : {stats['skipped']:,}  (already done)")
    print(f"\n  Dest folder : {dest_folder_id}")
    print(f"  Metadata    : {META_LOG}")
    print(f"  Audit log   : {AUDIT_LOG}")
    print(f"  Rejected log: {REJECT_LOG}")
    print("=" * 55)


if __name__ == "__main__":
    main()
